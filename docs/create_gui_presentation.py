"""
Script to capture GUI screenshots and generate a presentation with visual explanations
"""

import tkinter as tk
from PIL import Image, ImageDraw, ImageFont, ImageGrab
import numpy as np
import time
import os
import sys
import threading

# Add parent directory to path
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

def create_mock_gui_images():
    """Create mock GUI images for the presentation"""

    docs_dir = os.path.dirname(os.path.abspath(__file__))
    screenshots_dir = os.path.join(docs_dir, "screenshots")
    os.makedirs(screenshots_dir, exist_ok=True)

    # Colors matching the app theme
    BG_PRIMARY = (45, 45, 45)
    BG_SECONDARY = (54, 54, 54)
    BG_CARD = (58, 58, 58)
    CANVAS_BG = (30, 30, 30)
    TEXT_PRIMARY = (232, 232, 232)
    TEXT_SECONDARY = (176, 176, 176)
    ACCENT = (108, 155, 206)
    BORDER = (74, 74, 74)
    MARKER_COLOR = (255, 107, 107)

    def draw_rounded_rect(draw, xy, radius, fill, outline=None):
        """Draw a rounded rectangle"""
        x1, y1, x2, y2 = xy
        draw.rectangle([x1 + radius, y1, x2 - radius, y2], fill=fill)
        draw.rectangle([x1, y1 + radius, x2, y2 - radius], fill=fill)
        draw.ellipse([x1, y1, x1 + 2*radius, y1 + 2*radius], fill=fill)
        draw.ellipse([x2 - 2*radius, y1, x2, y1 + 2*radius], fill=fill)
        draw.ellipse([x1, y2 - 2*radius, x1 + 2*radius, y2], fill=fill)
        draw.ellipse([x2 - 2*radius, y2 - 2*radius, x2, y2], fill=fill)

    # === Screenshot 1: Main Interface Overview ===
    img1 = Image.new('RGB', (1200, 700), BG_PRIMARY)
    draw1 = ImageDraw.Draw(img1)

    # Control panel
    draw1.rectangle([0, 0, 1200, 50], fill=BG_SECONDARY)

    # Buttons
    buttons = ["Load Image 1", "Load Image 2", "|", "Normalize", "|", "Save", "|", "Zoom", "−", "100", "%", "+", "|", "Center", "Fit"]
    x_pos = 15
    for btn in buttons:
        if btn == "|":
            draw1.line([x_pos, 10, x_pos, 40], fill=BORDER, width=1)
            x_pos += 15
        elif btn in ["−", "+", "100", "%"]:
            if btn == "100":
                draw1.rectangle([x_pos, 12, x_pos + 50, 38], fill=BG_PRIMARY, outline=BORDER)
                draw1.text((x_pos + 15, 18), btn, fill=TEXT_PRIMARY)
                x_pos += 55
            elif btn == "%":
                draw1.text((x_pos, 18), btn, fill=TEXT_SECONDARY)
                x_pos += 20
            else:
                draw1.rectangle([x_pos, 12, x_pos + 25, 38], fill=BG_CARD)
                draw1.text((x_pos + 8, 15), btn, fill=TEXT_PRIMARY)
                x_pos += 30
        elif btn == "Zoom":
            draw1.text((x_pos, 18), btn, fill=TEXT_SECONDARY)
            x_pos += 45
        elif btn == "Normalize":
            draw1.rectangle([x_pos, 12, x_pos + 85, 38], fill=ACCENT)
            draw1.text((x_pos + 10, 18), btn, fill=(255, 255, 255))
            x_pos += 95
        else:
            draw1.rectangle([x_pos, 12, x_pos + 90, 38], fill=BG_CARD)
            draw1.text((x_pos + 10, 18), btn, fill=TEXT_PRIMARY)
            x_pos += 100

    # Three panels
    panel_width = 380
    panel_height = 530
    panel_y = 60
    panel_titles = ["Image 1", "Image 2", "Difference"]

    for i, title in enumerate(panel_titles):
        x = 10 + i * 395

        # Panel frame
        draw1.rectangle([x, panel_y, x + panel_width, panel_y + panel_height], fill=BG_CARD, outline=BORDER)

        # Panel title
        draw1.text((x + 10, panel_y + 5), title, fill=TEXT_PRIMARY)

        # Info bar
        draw1.text((x + 10, panel_y + 25), "256 × 256", fill=TEXT_SECONDARY)
        draw1.text((x + panel_width - 100, panel_y + 25), "(78,78) 100×100", fill=TEXT_SECONDARY)

        # Canvas area
        canvas_y = panel_y + 45
        canvas_h = panel_height - 55
        draw1.rectangle([x + 5, canvas_y, x + panel_width - 5, canvas_y + canvas_h], fill=CANVAS_BG)

        # Draw sample image content
        if i == 0:  # Image 1 - gradient
            for row in range(20):
                for col in range(20):
                    gray = int(50 + row * 8 + col * 2)
                    px = x + 80 + col * 11
                    py = canvas_y + 100 + row * 11
                    draw1.rectangle([px, py, px + 10, py + 10], fill=(gray, gray, gray))
        elif i == 1:  # Image 2 - similar gradient with slight diff
            for row in range(20):
                for col in range(20):
                    gray = int(55 + row * 8 + col * 2)
                    px = x + 80 + col * 11
                    py = canvas_y + 100 + row * 11
                    draw1.rectangle([px, py, px + 10, py + 10], fill=(gray, gray, gray))
        else:  # Difference - green/red
            for row in range(20):
                for col in range(20):
                    if row < 10:
                        color = (0, 100 + row * 10, 0)  # Green
                    else:
                        color = (100 + (row-10) * 10, 0, 0)  # Red
                    px = x + 80 + col * 11
                    py = canvas_y + 100 + row * 11
                    draw1.rectangle([px, py, px + 10, py + 10], fill=color)

    # Pixel info bar
    draw1.rectangle([0, 600, 1200, 630], fill=BG_SECONDARY)
    draw1.text((15, 608), "Pixel", fill=TEXT_SECONDARY)
    draw1.text((60, 608), "Position (128, 128)  |  Image1: 28456  |  Image2: 29012  |  Diff: -556", fill=TEXT_PRIMARY)

    # Status bar
    draw1.rectangle([0, 630, 1200, 660], fill=BG_PRIMARY)
    draw1.text((15, 638), "Ready", fill=TEXT_SECONDARY)
    draw1.text((1100, 638), "(128, 128)", fill=TEXT_SECONDARY)

    img1.save(os.path.join(screenshots_dir, "01_main_interface.png"))
    print("Created: 01_main_interface.png")

    # === Screenshot 2: Control Panel Detail ===
    img2 = Image.new('RGB', (1200, 200), BG_PRIMARY)
    draw2 = ImageDraw.Draw(img2)

    # Title
    draw2.text((20, 10), "Control Panel Components", fill=ACCENT)
    draw2.line([20, 35, 1180, 35], fill=ACCENT, width=2)

    # Control panel
    draw2.rectangle([20, 50, 1180, 100], fill=BG_SECONDARY)

    # Buttons with annotations
    annotations = [
        (30, "Load Image 1", "Load first image"),
        (140, "Load Image 2", "Load second image"),
        (270, "Normalize", "Toggle normalization"),
        (380, "Save", "Save images"),
        (470, "Zoom −/+", "Adjust zoom level"),
        (600, "Center", "Go to center"),
        (700, "Fit", "Fit to window"),
    ]

    for x, label, desc in annotations:
        # Button
        if label == "Normalize":
            draw2.rectangle([x, 60, x + 85, 90], fill=ACCENT)
            draw2.text((x + 5, 68), label, fill=(255, 255, 255))
        else:
            draw2.rectangle([x, 60, x + 90, 90], fill=BG_CARD)
            draw2.text((x + 5, 68), label, fill=TEXT_PRIMARY)

        # Arrow and description
        draw2.line([x + 45, 95, x + 45, 130], fill=MARKER_COLOR, width=2)
        draw2.polygon([(x + 40, 130), (x + 50, 130), (x + 45, 140)], fill=MARKER_COLOR)
        draw2.text((x, 150), desc, fill=TEXT_SECONDARY)

    img2.save(os.path.join(screenshots_dir, "02_control_panel.png"))
    print("Created: 02_control_panel.png")

    # === Screenshot 3: Image Panels Detail ===
    img3 = Image.new('RGB', (1200, 500), BG_PRIMARY)
    draw3 = ImageDraw.Draw(img3)

    # Title
    draw3.text((20, 10), "Three-Panel View", fill=ACCENT)
    draw3.line([20, 35, 1180, 35], fill=ACCENT, width=2)

    panel_info = [
        ("Image 1", "Original image (before calibration)", (100, 180, 100)),
        ("Image 2", "Comparison image (after calibration)", (100, 100, 180)),
        ("Difference", "Pixel difference visualization", None),
    ]

    for i, (title, desc, color) in enumerate(panel_info):
        x = 30 + i * 390

        # Panel
        draw3.rectangle([x, 50, x + 360, 400], fill=BG_CARD, outline=BORDER)
        draw3.text((x + 10, 55), title, fill=TEXT_PRIMARY)
        draw3.text((x + 10, 75), "256 × 256", fill=TEXT_SECONDARY)

        # Canvas
        draw3.rectangle([x + 10, 95, x + 350, 350], fill=CANVAS_BG)

        # Sample content
        if color:
            for row in range(12):
                for col in range(12):
                    c = tuple(min(255, v + row * 5 + col * 3) for v in color)
                    px = x + 50 + col * 23
                    py = 120 + row * 18
                    draw3.rectangle([px, py, px + 20, py + 15], fill=c)
        else:
            # Difference view
            for row in range(12):
                for col in range(12):
                    if (row + col) % 3 == 0:
                        c = (0, 150 + row * 5, 0)  # Green
                    elif (row + col) % 3 == 1:
                        c = (150 + col * 5, 0, 0)  # Red
                    else:
                        c = (30, 30, 30)  # No diff
                    px = x + 50 + col * 23
                    py = 120 + row * 18
                    draw3.rectangle([px, py, px + 20, py + 15], fill=c)

        # Description
        draw3.text((x + 10, 360), desc, fill=TEXT_SECONDARY)

        # Arrow pointing to panel
        arrow_x = x + 180
        draw3.line([arrow_x, 410, arrow_x, 440], fill=MARKER_COLOR, width=2)
        draw3.polygon([(arrow_x - 5, 410), (arrow_x + 5, 410), (arrow_x, 400)], fill=MARKER_COLOR)

    # Legend for difference
    draw3.text((850, 420), "Legend:", fill=TEXT_PRIMARY)
    draw3.rectangle([920, 418, 940, 432], fill=(0, 180, 0))
    draw3.text((945, 420), "Image1 > Image2", fill=TEXT_SECONDARY)
    draw3.rectangle([1070, 418, 1090, 432], fill=(180, 0, 0))
    draw3.text((1095, 420), "Image1 < Image2", fill=TEXT_SECONDARY)

    img3.save(os.path.join(screenshots_dir, "03_panel_view.png"))
    print("Created: 03_panel_view.png")

    # === Screenshot 4: Zoom Feature ===
    img4 = Image.new('RGB', (1200, 450), BG_PRIMARY)
    draw4 = ImageDraw.Draw(img4)

    draw4.text((20, 10), "Zoom Feature - Synchronized Across All Panels", fill=ACCENT)
    draw4.line([20, 35, 1180, 35], fill=ACCENT, width=2)

    # Show three zoom levels
    zoom_levels = [("25%", 3), ("100%", 1), ("400%", 0.25)]

    for i, (zoom, scale) in enumerate(zoom_levels):
        x = 30 + i * 390

        draw4.text((x + 150, 50), zoom, fill=TEXT_PRIMARY)
        draw4.rectangle([x, 80, x + 360, 380], fill=CANVAS_BG, outline=BORDER)

        # Draw grid representing pixels
        pixel_size = int(15 / scale)
        if pixel_size < 3:
            pixel_size = 3
        if pixel_size > 80:
            pixel_size = 80

        for row in range(int(280 / pixel_size)):
            for col in range(int(340 / pixel_size)):
                gray = 80 + (row * 7 + col * 5) % 100
                px = x + 10 + col * pixel_size
                py = 90 + row * pixel_size
                if px < x + 350 and py < 370:
                    draw4.rectangle([px, py, px + pixel_size - 1, py + pixel_size - 1],
                                   fill=(gray, gray, gray))

        # Zoom indicator
        draw4.text((x + 10, 385), f"Zoom: {zoom}", fill=TEXT_SECONDARY)

    # Instruction
    draw4.text((400, 420), "Use mouse wheel to zoom in/out - all three panels zoom together", fill=TEXT_PRIMARY)

    img4.save(os.path.join(screenshots_dir, "04_zoom_feature.png"))
    print("Created: 04_zoom_feature.png")

    # === Screenshot 5: Pixel Info Display ===
    img5 = Image.new('RGB', (1200, 400), BG_PRIMARY)
    draw5 = ImageDraw.Draw(img5)

    draw5.text((20, 10), "Pixel Value Display - Click to View 16-bit Values", fill=ACCENT)
    draw5.line([20, 35, 1180, 35], fill=ACCENT, width=2)

    # Show panel with marker
    draw5.rectangle([50, 60, 450, 300], fill=CANVAS_BG, outline=BORDER)
    draw5.text((220, 45), "Image 1", fill=TEXT_PRIMARY)

    # Draw some pixels
    for row in range(10):
        for col in range(10):
            gray = 100 + row * 10 + col * 5
            px = 100 + col * 30
            py = 90 + row * 20
            draw5.rectangle([px, py, px + 28, py + 18], fill=(gray, gray, gray))

    # Marker at clicked position
    marker_x, marker_y = 250, 180
    draw5.ellipse([marker_x - 8, marker_y - 8, marker_x + 8, marker_y + 8], outline=MARKER_COLOR, width=2)
    draw5.line([marker_x - 12, marker_y, marker_x + 12, marker_y], fill=MARKER_COLOR, width=1)
    draw5.line([marker_x, marker_y - 12, marker_x, marker_y + 12], fill=MARKER_COLOR, width=1)

    # Arrow from marker to info
    draw5.line([marker_x + 15, marker_y, 500, marker_y], fill=MARKER_COLOR, width=2)
    draw5.line([500, marker_y, 500, 320], fill=MARKER_COLOR, width=2)

    # Pixel info panel
    draw5.rectangle([480, 320, 1150, 380], fill=BG_SECONDARY, outline=ACCENT)
    draw5.text((500, 330), "Pixel Information (16-bit values)", fill=ACCENT)
    draw5.text((500, 355), "Position: (128, 95)  |  Image1: 28456  |  Image2: 29012  |  Diff: -556", fill=TEXT_PRIMARY)

    # Explanation
    draw5.rectangle([520, 60, 1150, 180], fill=BG_CARD)
    draw5.text((540, 70), "Click on any panel to:", fill=TEXT_PRIMARY)
    draw5.text((540, 95), "• Show cross marker at clicked position", fill=TEXT_SECONDARY)
    draw5.text((540, 115), "• Display actual 16-bit pixel values", fill=TEXT_SECONDARY)
    draw5.text((540, 135), "• Show difference value (Image1 - Image2)", fill=TEXT_SECONDARY)
    draw5.text((540, 155), "• Markers appear on all three panels", fill=TEXT_SECONDARY)

    img5.save(os.path.join(screenshots_dir, "05_pixel_info.png"))
    print("Created: 05_pixel_info.png")

    # === Screenshot 6: Normalization Feature ===
    img6 = Image.new('RGB', (1200, 450), BG_PRIMARY)
    draw6 = ImageDraw.Draw(img6)

    draw6.text((20, 10), "Normalization Feature - Standardize Brightness for Comparison", fill=ACCENT)
    draw6.line([20, 35, 1180, 35], fill=ACCENT, width=2)

    # Before normalization
    draw6.text((200, 50), "Before Normalization", fill=TEXT_PRIMARY)
    draw6.rectangle([50, 80, 350, 330], fill=CANVAS_BG, outline=BORDER)
    draw6.text((60, 85), "Image 1 (bright)", fill=TEXT_SECONDARY)
    for row in range(10):
        for col in range(10):
            gray = 150 + row * 5 + col * 3
            draw6.rectangle([70 + col * 26, 110 + row * 20, 94 + col * 26, 128 + row * 20],
                           fill=(gray, gray, gray))

    draw6.rectangle([380, 80, 680, 330], fill=CANVAS_BG, outline=BORDER)
    draw6.text((390, 85), "Image 2 (dark)", fill=TEXT_SECONDARY)
    for row in range(10):
        for col in range(10):
            gray = 50 + row * 5 + col * 3
            draw6.rectangle([400 + col * 26, 110 + row * 20, 424 + col * 26, 128 + row * 20],
                           fill=(gray, gray, gray))

    # Arrow
    draw6.line([700, 200, 780, 200], fill=ACCENT, width=3)
    draw6.polygon([(780, 193), (780, 207), (800, 200)], fill=ACCENT)
    draw6.text((720, 220), "Normalize", fill=ACCENT)

    # After normalization
    draw6.text((950, 50), "After Normalization", fill=TEXT_PRIMARY)
    draw6.rectangle([820, 80, 1120, 330], fill=CANVAS_BG, outline=BORDER)
    draw6.text((830, 85), "Both images normalized", fill=TEXT_SECONDARY)
    for row in range(10):
        for col in range(10):
            gray = 117 + row * 3 + col * 2  # Similar brightness
            draw6.rectangle([840 + col * 26, 110 + row * 20, 864 + col * 26, 128 + row * 20],
                           fill=(gray, gray, gray))

    # Explanation
    draw6.rectangle([50, 350, 1150, 430], fill=BG_CARD)
    draw6.text((70, 360), "Normalization Method:", fill=ACCENT)
    draw6.text((70, 385), "• Calculates average of center 100×100 pixel region", fill=TEXT_SECONDARY)
    draw6.text((70, 405), "• Scales all pixel values so center average = 30000 (16-bit target)", fill=TEXT_SECONDARY)
    draw6.text((600, 385), "• Toggle: Click 'Normalize' to apply, click 'Restore' to revert", fill=TEXT_SECONDARY)
    draw6.text((600, 405), "• Both raw 16-bit data and display image are normalized", fill=TEXT_SECONDARY)

    img6.save(os.path.join(screenshots_dir, "06_normalization.png"))
    print("Created: 06_normalization.png")

    # === Screenshot 7: Difference Visualization ===
    img7 = Image.new('RGB', (1200, 500), BG_PRIMARY)
    draw7 = ImageDraw.Draw(img7)

    draw7.text((20, 10), "Difference Visualization - Color-Coded Pixel Differences", fill=ACCENT)
    draw7.line([20, 35, 1180, 35], fill=ACCENT, width=2)

    # Image 1
    draw7.rectangle([50, 60, 300, 310], fill=CANVAS_BG, outline=BORDER)
    draw7.text((140, 45), "Image 1", fill=TEXT_PRIMARY)
    for row in range(10):
        for col in range(10):
            gray = 100 + row * 8 + col * 4
            draw7.rectangle([70 + col * 22, 80 + row * 22, 90 + col * 22, 100 + row * 22],
                           fill=(gray, gray, gray))

    # Minus sign
    draw7.text((330, 170), "−", fill=TEXT_PRIMARY)

    # Image 2
    draw7.rectangle([370, 60, 620, 310], fill=CANVAS_BG, outline=BORDER)
    draw7.text((460, 45), "Image 2", fill=TEXT_PRIMARY)
    for row in range(10):
        for col in range(10):
            # Slightly different values
            if row < 5:
                gray = 95 + row * 8 + col * 4  # Darker than Image 1 (positive diff)
            else:
                gray = 110 + row * 8 + col * 4  # Brighter than Image 1 (negative diff)
            draw7.rectangle([390 + col * 22, 80 + row * 22, 410 + col * 22, 100 + row * 22],
                           fill=(gray, gray, gray))

    # Equals sign
    draw7.text((650, 170), "=", fill=TEXT_PRIMARY)

    # Difference
    draw7.rectangle([700, 60, 950, 310], fill=CANVAS_BG, outline=BORDER)
    draw7.text((780, 45), "Difference", fill=TEXT_PRIMARY)
    for row in range(10):
        for col in range(10):
            if row < 5:
                # Positive diff: Image1 > Image2 → GREEN
                intensity = 80 + row * 20 + col * 5
                color = (0, intensity, 0)
            else:
                # Negative diff: Image1 < Image2 → RED
                intensity = 80 + (row - 5) * 20 + col * 5
                color = (intensity, 0, 0)
            draw7.rectangle([720 + col * 22, 80 + row * 22, 740 + col * 22, 100 + row * 22],
                           fill=color)

    # Legend
    draw7.rectangle([980, 60, 1170, 200], fill=BG_CARD, outline=BORDER)
    draw7.text((1000, 70), "Color Legend", fill=ACCENT)

    draw7.rectangle([1000, 100, 1030, 120], fill=(0, 180, 0))
    draw7.text((1040, 102), "GREEN", fill=(0, 180, 0))
    draw7.text((1000, 125), "Image1 > Image2", fill=TEXT_SECONDARY)
    draw7.text((1000, 145), "(Positive difference)", fill=TEXT_SECONDARY)

    draw7.rectangle([1000, 170, 1030, 190], fill=(180, 0, 0))
    draw7.text((1040, 172), "RED", fill=(180, 0, 0))
    draw7.text((1000, 195), "Image1 < Image2", fill=TEXT_SECONDARY)
    draw7.text((1000, 215), "(Negative difference)", fill=TEXT_SECONDARY)

    # Formula
    draw7.rectangle([50, 330, 950, 480], fill=BG_CARD)
    draw7.text((70, 345), "Calculation Formula:", fill=ACCENT)
    draw7.text((70, 375), "diff = Image1_pixel - Image2_pixel    (signed difference, preserves direction)", fill=TEXT_PRIMARY)
    draw7.text((70, 405), "• Positive diff → Display in GREEN channel (intensity proportional to diff magnitude)", fill=TEXT_SECONDARY)
    draw7.text((70, 430), "• Negative diff → Display in RED channel (intensity proportional to |diff| magnitude)", fill=TEXT_SECONDARY)
    draw7.text((70, 455), "• Zero diff → Black (no visible difference)", fill=TEXT_SECONDARY)

    img7.save(os.path.join(screenshots_dir, "07_difference_view.png"))
    print("Created: 07_difference_view.png")

    # === Screenshot 8: Save Feature ===
    img8 = Image.new('RGB', (1200, 400), BG_PRIMARY)
    draw8 = ImageDraw.Draw(img8)

    draw8.text((20, 10), "Save Feature - Export Normalized Images", fill=ACCENT)
    draw8.line([20, 35, 1180, 35], fill=ACCENT, width=2)

    # Save button
    draw8.rectangle([50, 60, 150, 100], fill=BG_CARD, outline=ACCENT)
    draw8.text((75, 72), "Save", fill=TEXT_PRIMARY)

    # Arrow
    draw8.line([160, 80, 250, 80], fill=ACCENT, width=2)
    draw8.polygon([(250, 73), (250, 87), (270, 80)], fill=ACCENT)

    # Output files
    draw8.rectangle([300, 50, 1150, 380], fill=BG_CARD)
    draw8.text((320, 60), "Output Files (16-bit TIFF format):", fill=ACCENT)

    files = [
        ("image1_normalized.tif", "Image 1 × 16 (scaled to full 16-bit range)"),
        ("image2_normalized.tif", "Image 2 × 16 (scaled to full 16-bit range)"),
        ("diff_normalized.tif", "Difference × 16 + 30000 (centered at 30000)"),
        ("diff_visualization.png", "Color-coded difference image (RGB)"),
    ]

    for i, (filename, desc) in enumerate(files):
        y = 100 + i * 65
        draw8.rectangle([330, y, 370, y + 40], fill=CANVAS_BG)
        draw8.text((340, y + 12), "📄", fill=TEXT_PRIMARY)
        draw8.text((390, y + 5), filename, fill=TEXT_PRIMARY)
        draw8.text((390, y + 25), desc, fill=TEXT_SECONDARY)

    img8.save(os.path.join(screenshots_dir, "08_save_feature.png"))
    print("Created: 08_save_feature.png")

    print(f"\nAll screenshots saved to: {screenshots_dir}")
    return screenshots_dir


def create_gui_presentation(screenshots_dir):
    """Create presentation with GUI screenshots"""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def add_slide_with_image(title, image_path, caption=""):
        """Add a slide with title and image"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(0x2D, 0x2D, 0x2D)

        # Line under title
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.85), Inches(12.333), Inches(0.03))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(0x6C, 0x9B, 0xCE)
        line.line.fill.background()

        # Image
        if os.path.exists(image_path):
            img = slide.shapes.add_picture(image_path, Inches(0.5), Inches(1), Inches(12.333), Inches(5.5))

        # Caption
        if caption:
            caption_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.7), Inches(12), Inches(0.5))
            caption_frame = caption_box.text_frame
            caption_para = caption_frame.paragraphs[0]
            caption_para.text = caption
            caption_para.font.size = Pt(14)
            caption_para.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
            caption_para.alignment = PP_ALIGN.CENTER

        return slide

    def add_title_slide(title, subtitle):
        """Add a title slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Background accent bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.8), Inches(13.333), Inches(1.8))
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(0x6C, 0x9B, 0xCE)
        bar.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.9), Inches(12.333), Inches(1))
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(48)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        title_para.alignment = PP_ALIGN.CENTER

        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(12.333), Inches(0.8))
        subtitle_frame = subtitle_box.text_frame
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.text = subtitle
        subtitle_para.font.size = Pt(24)
        subtitle_para.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
        subtitle_para.alignment = PP_ALIGN.CENTER

        return slide

    # === Slides ===

    # Slide 1: Title
    add_title_slide(
        "Image Compare Tool",
        "GUI Visual Guide\nByung Geun (BG) Jun | March 2026"
    )

    # Slide 2: Main Interface
    add_slide_with_image(
        "Main Interface Overview",
        os.path.join(screenshots_dir, "01_main_interface.png"),
        "Three-panel view with Control Panel, Image Panels, and Status Bars"
    )

    # Slide 3: Control Panel
    add_slide_with_image(
        "Control Panel Components",
        os.path.join(screenshots_dir, "02_control_panel.png"),
        "All controls for image loading, normalization, saving, and navigation"
    )

    # Slide 4: Panel View
    add_slide_with_image(
        "Three-Panel Image View",
        os.path.join(screenshots_dir, "03_panel_view.png"),
        "Left: Image 1 | Center: Image 2 | Right: Difference visualization"
    )

    # Slide 5: Zoom Feature
    add_slide_with_image(
        "Synchronized Zoom Feature",
        os.path.join(screenshots_dir, "04_zoom_feature.png"),
        "All three panels zoom together - use mouse wheel or +/- buttons"
    )

    # Slide 6: Pixel Info
    add_slide_with_image(
        "Pixel Value Display",
        os.path.join(screenshots_dir, "05_pixel_info.png"),
        "Click on any panel to view actual 16-bit pixel values"
    )

    # Slide 7: Normalization
    add_slide_with_image(
        "Image Normalization",
        os.path.join(screenshots_dir, "06_normalization.png"),
        "Standardize brightness based on center 100×100 region average"
    )

    # Slide 8: Difference View
    add_slide_with_image(
        "Difference Visualization",
        os.path.join(screenshots_dir, "07_difference_view.png"),
        "Color-coded difference: GREEN = Image1 > Image2, RED = Image1 < Image2"
    )

    # Slide 9: Save Feature
    add_slide_with_image(
        "Save Feature",
        os.path.join(screenshots_dir, "08_save_feature.png"),
        "Export normalized images as 16-bit TIFF files"
    )

    # Slide 10: Thank You
    add_title_slide(
        "Thank You",
        "Questions & Discussion\n\nContact: byungpaul@meta.com"
    )

    # Save
    output_path = os.path.join(os.path.dirname(screenshots_dir), "Image_Compare_Tool_GUI_Guide.pptx")
    prs.save(output_path)
    print(f"\nPresentation saved to: {output_path}")
    return output_path


if __name__ == "__main__":
    print("Creating GUI screenshots...")
    screenshots_dir = create_mock_gui_images()

    print("\nCreating presentation with GUI images...")
    create_gui_presentation(screenshots_dir)

    print("\nDone!")
