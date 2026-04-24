"""
Script to generate an English PowerPoint presentation for Image Compare Tool
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0x2D, 0x2D, 0x2D)
    title_para.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.text = subtitle
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    subtitle_para.alignment = PP_ALIGN.CENTER

    return slide

def add_section_slide(prs, title):
    """Add a section divider slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add background shape
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.5), Inches(10), Inches(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x6C, 0x9B, 0xCE)
    shape.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    title_para.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, bullet_points):
    """Add a content slide with bullet points"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0x2D, 0x2D, 0x2D)

    # Add line under title
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(9), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0x6C, 0x9B, 0xCE)
    line.line.fill.background()

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(5.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    for i, point in enumerate(bullet_points):
        if i == 0:
            para = content_frame.paragraphs[0]
        else:
            para = content_frame.add_paragraph()

        para.text = f"• {point}"
        para.font.size = Pt(20)
        para.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        para.space_before = Pt(12)
        para.space_after = Pt(6)

    return slide

def add_two_column_slide(prs, title, left_content, right_content):
    """Add a two-column content slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0x2D, 0x2D, 0x2D)

    # Line under title
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(9), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0x6C, 0x9B, 0xCE)
    line.line.fill.background()

    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(4.3), Inches(5.5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True

    for i, (label, value) in enumerate(left_content):
        if i == 0:
            para = left_frame.paragraphs[0]
        else:
            para = left_frame.add_paragraph()

        para.text = f"• {label}: {value}"
        para.font.size = Pt(18)
        para.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        para.space_before = Pt(8)

    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.4), Inches(4.3), Inches(5.5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True

    for i, (label, value) in enumerate(right_content):
        if i == 0:
            para = right_frame.paragraphs[0]
        else:
            para = right_frame.add_paragraph()

        para.text = f"• {label}: {value}"
        para.font.size = Pt(18)
        para.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        para.space_before = Pt(8)

    return slide

def add_code_slide(prs, title, code_text, description=""):
    """Add a slide with code block"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0x2D, 0x2D, 0x2D)

    # Description
    if description:
        desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(0.6))
        desc_frame = desc_box.text_frame
        desc_para = desc_frame.paragraphs[0]
        desc_para.text = description
        desc_para.font.size = Pt(16)
        desc_para.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    # Code block background
    code_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.8), Inches(9), Inches(4.5))
    code_bg.fill.solid()
    code_bg.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x1E)
    code_bg.line.fill.background()

    # Code text
    code_box = slide.shapes.add_textbox(Inches(0.7), Inches(2), Inches(8.6), Inches(4.2))
    code_frame = code_box.text_frame
    code_frame.word_wrap = True

    for i, line in enumerate(code_text.split('\n')):
        if i == 0:
            para = code_frame.paragraphs[0]
        else:
            para = code_frame.add_paragraph()

        para.text = line
        para.font.size = Pt(12)
        para.font.name = 'Consolas'
        para.font.color.rgb = RGBColor(0xE8, 0xE8, 0xE8)

    return slide

def create_presentation():
    """Create the main presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # === Slide 1: Title ===
    add_title_slide(
        prs,
        "Image Compare Tool",
        "Intermediate Development Report\nByung Geun (BG) Jun | March 2026"
    )

    # === Slide 2: Agenda ===
    add_content_slide(prs, "Agenda", [
        "Project Overview",
        "Key Features Implemented",
        "Technical Architecture",
        "16-bit Image Processing",
        "Normalization Feature",
        "UI/UX Design",
        "Bug Fixes & Improvements",
        "Current Status & Next Steps"
    ])

    # === Slide 3: Section - Overview ===
    add_section_slide(prs, "Project Overview")

    # === Slide 4: Project Overview ===
    add_content_slide(prs, "Project Overview", [
        "Purpose: Compare two images and visualize pixel differences",
        "Target Users: Display engineers, Image processing specialists",
        "Primary Use Case: Calibration before/after comparison",
        "Technology Stack: Python, Tkinter, PIL/Pillow, NumPy",
        "Platform: Windows (standalone executable available)"
    ])

    # === Slide 5: Section - Features ===
    add_section_slide(prs, "Key Features")

    # === Slide 6: Core Features ===
    add_two_column_slide(
        prs,
        "Core Features Implemented",
        [
            ("Dual Image Loading", "Independent load for each image"),
            ("Three-Panel View", "Image1, Image2, Difference"),
            ("16-bit TIFF Support", "Full 16-bit data preservation"),
            ("Pixel Value Display", "Click to see actual values"),
            ("Synchronized Zoom", "All panels zoom together"),
        ],
        [
            ("Pan/Drag Navigation", "Click and drag to move"),
            ("Normalization", "Center 100×100 region based"),
            ("Difference Visualization", "Color-coded heatmap"),
            ("Image Save", "Export normalized images"),
            ("Marker Display", "Click position indicator"),
        ]
    )

    # === Slide 7: Difference Visualization ===
    add_content_slide(prs, "Difference Visualization", [
        "Calculates actual pixel difference: Image1 - Image2",
        "Color-coded display based on signed difference values",
        "GREEN channel: Positive difference (Image1 > Image2)",
        "RED channel: Negative difference (Image1 < Image2)",
        "Intensity proportional to difference magnitude",
        "Preserves sign information for analysis"
    ])

    # === Slide 8: Section - Technical ===
    add_section_slide(prs, "Technical Architecture")

    # === Slide 9: Architecture ===
    add_content_slide(prs, "Software Architecture", [
        "Modular class-based design",
        "Theme class: Clean minimal dark theme with consistent styling",
        "ImagePanel class: Individual panel management with markers",
        "ImageCompareApp class: Main application controller",
        "Event-driven architecture for user interactions",
        "NumPy-based image processing for performance"
    ])

    # === Slide 10: 16-bit Processing ===
    add_code_slide(
        prs,
        "16-bit Image Processing",
        """# Load 16-bit TIFF image
if image.mode == 'I;16':
    arr = np.array(image, dtype=np.uint16)

    # Store original 16-bit data
    self.image1_raw = arr.copy()

    # Convert to 8-bit for display using fixed range
    # Key fix: Use 65535 instead of arr.max()
    arr_display = (arr.astype(np.float32) / 65535.0 * 255)
    arr_display = arr_display.astype(np.uint8)

    image = Image.fromarray(arr_display)""",
        "Preserves original 16-bit data while displaying normalized 8-bit version"
    )

    # === Slide 11: Section - Normalization ===
    add_section_slide(prs, "Normalization Feature")

    # === Slide 12: Normalization Details ===
    add_content_slide(prs, "Image Normalization", [
        "Purpose: Standardize brightness for fair comparison",
        "Method: Center 100×100 region average → Target value (30000)",
        "Target value based on 16-bit range (0-65535)",
        "Toggle functionality: Normalize / Restore original",
        "Both raw data and display image are normalized",
        "Difference image recalculated after normalization"
    ])

    # === Slide 13: Normalization Code ===
    add_code_slide(
        prs,
        "Normalization Algorithm",
        """# Calculate center 100x100 region average
h, w = arr.shape[:2]
center_x, center_y = w // 2, h // 2
x1, y1 = max(0, center_x - 50), max(0, center_y - 50)
x2, y2 = min(w, center_x + 50), min(h, center_y + 50)

center_region = arr[y1:y2, x1:x2]
current_avg = np.mean(center_region)

# Calculate normalization factor
NORMALIZATION_TARGET = 30000  # 16-bit based
factor = NORMALIZATION_TARGET / (current_avg * 256)

# Apply normalization
normalized_arr = np.clip(arr * factor, 0, 255).astype(np.uint8)""",
        "Normalizes based on center region average to target value"
    )

    # === Slide 14: Section - UI/UX ===
    add_section_slide(prs, "UI/UX Design")

    # === Slide 15: Theme Design ===
    add_two_column_slide(
        prs,
        "Clean Minimal Theme",
        [
            ("Background", "#2D2D2D (Dark Gray)"),
            ("Cards", "#3A3A3A (Slightly lighter)"),
            ("Canvas", "#1E1E1E (Near black)"),
            ("Text Primary", "#E8E8E8 (Off-white)"),
            ("Text Secondary", "#B0B0B0 (Gray)"),
        ],
        [
            ("Accent Color", "#6C9BCE (Soft Blue)"),
            ("Borders", "#4A4A4A (Subtle)"),
            ("Marker", "#FF6B6B (Coral Red)"),
            ("Font", "Segoe UI / Consolas"),
            ("Design", "Flat, minimal, professional"),
        ]
    )

    # === Slide 16: UI Layout ===
    add_content_slide(prs, "Interface Layout", [
        "Control Panel (Top): Load, Normalize, Save, Zoom controls",
        "Three-Panel View (Center): Image1, Image2, Difference",
        "Each panel shows: Image size, View position, Current region",
        "Status Bar (Bottom): Ready state, Current position",
        "Pixel Info Bar: Click to display 16-bit values",
        "Cross marker on click position across all panels"
    ])

    # === Slide 17: Section - Bug Fixes ===
    add_section_slide(prs, "Bug Fixes & Improvements")

    # === Slide 18: Bug Fixes ===
    add_content_slide(prs, "Recent Bug Fixes", [
        "Fixed: Image 2 appearing too dark after loading",
        "Cause: Using arr.max() for normalization caused inconsistent scaling",
        "Solution: Use fixed 65535 (16-bit max) for consistent brightness",
        "Fixed: Restore function also had same brightness issue",
        "All Korean comments translated to English",
        "Code documentation improved for maintainability"
    ])

    # === Slide 19: Before/After Fix ===
    add_code_slide(
        prs,
        "Brightness Fix - Before vs After",
        """# BEFORE (Bug): Each image scaled differently
arr_display = (arr.astype(np.float32) / arr.max() * 255)
# Image1 max=30000 → pixels scaled by 30000
# Image2 max=60000 → pixels scaled by 60000
# Same pixel value appears different brightness!

# AFTER (Fixed): Consistent scaling for all images
arr_display = (arr.astype(np.float32) / 65535.0 * 255)
# Both images use same 16-bit full range
# Same pixel value = Same brightness
# Proper visual comparison possible""",
        "Fixed brightness inconsistency between loaded images"
    )

    # === Slide 20: Section - Status ===
    add_section_slide(prs, "Current Status")

    # === Slide 21: Current Status ===
    add_two_column_slide(
        prs,
        "Development Status",
        [
            ("Core Features", "✓ Complete"),
            ("16-bit Support", "✓ Complete"),
            ("Normalization", "✓ Complete"),
            ("Image Save", "✓ Complete"),
            ("UI Theme", "✓ Complete"),
        ],
        [
            ("Pixel Info Display", "✓ Complete"),
            ("Marker System", "✓ Complete"),
            ("Bug Fixes", "✓ Complete"),
            ("English Translation", "✓ Complete"),
            ("Documentation", "In Progress"),
        ]
    )

    # === Slide 22: Next Steps ===
    add_content_slide(prs, "Next Steps", [
        "Add ROI (Region of Interest) selection feature",
        "Implement histogram comparison view",
        "Add measurement tools (distance, angle)",
        "Support for batch image comparison",
        "Export comparison report (PDF/HTML)",
        "Performance optimization for large images",
        "Cross-platform support (Mac/Linux)"
    ])

    # === Slide 23: Thank You ===
    add_title_slide(
        prs,
        "Thank You",
        "Questions & Discussion\n\nContact: byungpaul@meta.com"
    )

    # Save the presentation
    output_path = os.path.join(os.path.dirname(__file__), "Image_Compare_Tool_Report.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    return output_path

if __name__ == "__main__":
    create_presentation()
